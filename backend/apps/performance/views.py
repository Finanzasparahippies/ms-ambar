from rest_framework import viewsets, status
from rest_framework.decorators import action
from rest_framework.response import Response
from rest_framework.permissions import AllowAny, IsAdminUser
from django.db.models import Avg, Max, Count
from .models import PerformanceMetric, ServerRequestLog
from .serializers import PerformanceMetricSerializer, ServerRequestLogSerializer

class PerformanceViewSet(viewsets.ViewSet):
    permission_classes = [AllowAny]

    @action(detail=False, methods=['post'], url_path='vitals')
    def report_vitals(self, request):
        """Endpoint for the frontend to report Web Vitals (supports single & batch array)."""
        data = request.data
        if isinstance(data, list):
            serializer = PerformanceMetricSerializer(data=data, many=True)
            if serializer.is_valid():
                serializer.save(user_agent=request.META.get('HTTP_USER_AGENT', ''))
                return Response(status=status.HTTP_201_CREATED)
            return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

        serializer = PerformanceMetricSerializer(data=data)
        if serializer.is_valid():
            serializer.save(user_agent=request.META.get('HTTP_USER_AGENT', ''))
            return Response(status=status.HTTP_201_CREATED)
        return Response(serializer.errors, status=status.HTTP_400_BAD_REQUEST)

    @action(detail=False, methods=['get'], url_path='summary', permission_classes=[IsAdminUser])
    def get_summary(self, request):
        """Returns a summary of performance metrics for the admin dashboard."""
        server_summary = ServerRequestLog.objects.aggregate(
            avg_response_time=Avg('response_time'),
            max_response_time=Max('response_time'),
            avg_queries=Avg('query_count'),
            total_requests=Count('id')
        )
        
        vitals_summary = PerformanceMetric.objects.values('name').annotate(
            avg_value=Avg('value'),
            count=Count('id')
        )

        slowest_endpoints = ServerRequestLog.objects.values('path').annotate(
            avg_time=Avg('response_time')
        ).order_by('-avg_time')[:10]

        return Response({
            'server': server_summary,
            'vitals': vitals_summary,
            'slowest_endpoints': slowest_endpoints
        })

    @action(detail=False, methods=['get'], url_path='logs', permission_classes=[IsAdminUser])
    def list_logs(self, request):
        """Lists available log files in the logs directory, including active and rotated files."""
        from django.conf import settings
        import os

        logs_dir = getattr(settings, 'LOGS_DIR', settings.BASE_DIR / 'logs')
        # Crear dinámicamente si no existe para evitar errores E/S
        if not os.path.exists(logs_dir):
            os.makedirs(logs_dir, exist_ok=True)

        allowed_bases = ('tickets.log', 'shop.log', 'events.log', 'users.log', 'blog.log', 'dashboard.log', 'gallery.log', 'music.log', 'production.log', 'staging.log', 'test.log', 'local.log')
        log_files = []

        # Obtener los archivos físicos en el directorio de logs que inicien con nuestras bases permitidas
        existing_files = {}
        try:
            for file_name in os.listdir(logs_dir):
                if file_name.startswith(allowed_bases):
                    file_path = os.path.join(logs_dir, file_name)
                    if os.path.isfile(file_path):
                        stat = os.stat(file_path)
                        existing_files[file_name] = {
                            'name': file_name,
                            'size': stat.st_size,
                            'modified': stat.st_mtime
                        }
        except Exception:
            # En caso de error de lectura, continuamos para no romper el admin
            pass

        # Asegurar que las bases principales siempre aparezcan listadas,
        # aunque no se hayan escrito aún físicamente en disco.
        for base_name in allowed_bases:
            if base_name in existing_files:
                log_files.append(existing_files[base_name])
            else:
                log_files.append({
                    'name': base_name,
                    'size': 0,
                    'modified': None
                })

            # También añadimos los archivos rotados que compartan esa misma base ordenados
            for file_name, file_info in sorted(existing_files.items()):
                if file_name != base_name and file_name.startswith(base_name):
                    log_files.append(file_info)

        return Response(log_files)

    @action(detail=False, methods=['get'], url_path='logs/download', permission_classes=[IsAdminUser])
    def download_log(self, request):
        """Downloads a specific log file, preventing path traversal attacks."""
        from django.conf import settings
        from django.http import FileResponse, Http404
        import os

        file_name = request.query_params.get('file', '').strip()
        allowed_bases = ('tickets.log', 'shop.log', 'events.log', 'users.log', 'blog.log', 'dashboard.log', 'gallery.log', 'music.log', 'production.log', 'staging.log', 'test.log', 'local.log')

        # Validación estricta del nombre del archivo solicitado
        if not file_name or not file_name.startswith(allowed_bases):
            return Response({'error': 'Archivo no permitido.'}, status=status.HTTP_400_BAD_REQUEST)

        # Mitigación estricta de Path Traversal
        if '/' in file_name or '\\' in file_name or '..' in file_name:
            return Response({'error': 'Nombre de archivo inválido.'}, status=status.HTTP_400_BAD_REQUEST)

        logs_dir = getattr(settings, 'LOGS_DIR', settings.BASE_DIR / 'logs')
        file_path = os.path.join(logs_dir, file_name)

        # Si no existe y es una de las bases principales, lo creamos vacío en caliente
        # para evitar fallos de archivo no encontrado al intentar descargarlo
        if not os.path.exists(file_path):
            if file_name in allowed_bases:
                try:
                    os.makedirs(logs_dir, exist_ok=True)
                    with open(file_path, 'w', encoding='utf-8') as f:
                        pass
                except Exception as e:
                    return Response({'error': f'No se pudo inicializar el archivo de log: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
            else:
                raise Http404("El archivo de log no existe.")

        try:
            # FileResponse con as_attachment=True para forzar la descarga directa del navegador con charset UTF-8
            response = FileResponse(
                open(file_path, 'rb'),
                as_attachment=True,
                filename=file_name,
                content_type='text/plain; charset=utf-8'
            )
            return response
        except Exception as e:
            return Response({'error': f'Error al descargar el archivo: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    @action(detail=False, methods=['post'], url_path='logs/purge', permission_classes=[IsAdminUser])
    def purge_log(self, request):
        """Purges/empties the specified log file in real time using secure fcntl file locking."""
        from django.conf import settings
        import os

        file_name = (request.data.get('file') or request.query_params.get('file') or '').strip()
        allowed_bases = ('tickets.log', 'shop.log', 'events.log', 'users.log', 'blog.log', 'dashboard.log', 'gallery.log', 'music.log', 'production.log', 'staging.log', 'test.log', 'local.log')

        if not file_name or not file_name.startswith(allowed_bases):
            return Response({'error': 'Archivo no permitido para purga.'}, status=status.HTTP_400_BAD_REQUEST)

        if '/' in file_name or '\\' in file_name or '..' in file_name:
            return Response({'error': 'Nombre de archivo inválido.'}, status=status.HTTP_400_BAD_REQUEST)

        logs_dir = getattr(settings, 'LOGS_DIR', settings.BASE_DIR / 'logs')
        file_path = os.path.join(logs_dir, file_name)

        if not os.path.exists(file_path):
            return Response({'error': f'El archivo {file_name} no existe en disco.'}, status=status.HTTP_404_NOT_FOUND)

        try:
            # Intentar bloqueo seguro mediante fcntl (Linux/Hetzner) antes de truncar
            with open(file_path, 'r+', encoding='utf-8') as f:
                try:
                    import fcntl
                    fcntl.flock(f, fcntl.LOCK_EX)
                except Exception:
                    pass
                f.truncate(0)
                f.seek(0)
                try:
                    import fcntl
                    fcntl.flock(f, fcntl.LOCK_UN)
                except Exception:
                    pass
            return Response({'message': f'Log {file_name} purgado exitosamente.', 'file': file_name}, status=status.HTTP_200_OK)
        except Exception as e:
            return Response({'error': f'No se pudo purgar el archivo: {str(e)}'}, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

    @action(detail=False, methods=['get'], url_path='export/pdf', permission_classes=[IsAdminUser])
    def export_pdf(self, request):
        """Generates an executive PDF report of system performance and logs audit."""
        import io
        from fpdf import FPDF
        from django.http import HttpResponse
        from datetime import datetime

        server_summary = ServerRequestLog.objects.aggregate(
            avg_response_time=Avg('response_time'),
            max_response_time=Max('response_time'),
            avg_queries=Avg('query_count'),
            total_requests=Count('id')
        )

        vitals = PerformanceMetric.objects.values('name').annotate(
            avg_value=Avg('value'),
            count=Count('id')
        )

        slowest = ServerRequestLog.objects.values('path').annotate(
            avg_time=Avg('response_time')
        ).order_by('-avg_time')[:10]

        pdf = FPDF()
        pdf.add_page()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Encabezado Corporativo
        pdf.set_fill_color(8, 12, 10)
        pdf.rect(0, 0, 210, 28, 'F')
        pdf.set_font('Helvetica', 'B', 14)
        pdf.set_text_color(245, 158, 11)
        pdf.cell(0, 10, 'NECTAR LABS | MS-AMBAR PERFORMANCE REPORT', 0, 1, 'C')
        pdf.set_font('Helvetica', '', 9)
        pdf.set_text_color(200, 200, 200)
        pdf.cell(0, 5, f'Fecha de Emision: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}', 0, 1, 'C')
        pdf.ln(12)

        # Seccion 1: Resumen de Servidor
        pdf.set_font('Helvetica', 'B', 12)
        pdf.set_text_color(10, 10, 10)
        pdf.cell(0, 8, '1. METRICAS CLAVE DE INFRAESTRUCTURA', 0, 1, 'L')
        pdf.set_draw_color(245, 158, 11)
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(4)

        pdf.set_font('Helvetica', '', 10)
        avg_resp = f"{server_summary.get('avg_response_time') or 0:.3f} s"
        max_resp = f"{server_summary.get('max_response_time') or 0:.3f} s"
        avg_queries = f"{server_summary.get('avg_queries') or 0:.1f}"
        total_reqs = str(server_summary.get('total_requests') or 0)

        pdf.cell(95, 8, f"Tiempo de Respuesta Promedio: {avg_resp}", 1, 0, 'L')
        pdf.cell(95, 8, f"Tiempo de Respuesta Maximo: {max_resp}", 1, 1, 'L')
        pdf.cell(95, 8, f"Consultas DB Promedio (Queries): {avg_queries}", 1, 0, 'L')
        pdf.cell(95, 8, f"Total de Solicitudes Registradas: {total_reqs}", 1, 1, 'L')
        pdf.ln(8)

        # Seccion 2: Core Web Vitals
        pdf.set_font('Helvetica', 'B', 12)
        pdf.cell(0, 8, '2. CORE WEB VITALS (PROMEDIO FRONTEND)', 0, 1, 'L')
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(4)

        if vitals:
            pdf.set_font('Helvetica', 'B', 9)
            pdf.set_fill_color(240, 240, 240)
            pdf.cell(60, 7, 'Metrica Vital', 1, 0, 'C', True)
            pdf.cell(65, 7, 'Valor Promedio (ms)', 1, 0, 'C', True)
            pdf.cell(65, 7, 'Total Muestras', 1, 1, 'C', True)
            pdf.set_font('Helvetica', '', 9)
            for v in vitals:
                pdf.cell(60, 7, str(v['name']), 1, 0, 'C')
                pdf.cell(65, 7, f"{v['avg_value']:.2f} ms", 1, 0, 'C')
                pdf.cell(65, 7, str(v['count']), 1, 1, 'C')
        else:
            pdf.set_font('Helvetica', 'I', 9)
            pdf.cell(0, 7, 'No hay muestras registradas de Web Vitals.', 0, 1, 'L')
        pdf.ln(8)

        # Seccion 3: Endpoints Mas Lentos
        pdf.set_font('Helvetica', 'B', 12)
        pdf.cell(0, 8, '3. LATENCIA CRITICA DE ENDPOINTS (TOP 10)', 0, 1, 'L')
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(4)

        if slowest:
            pdf.set_font('Helvetica', 'B', 9)
            pdf.set_fill_color(240, 240, 240)
            pdf.cell(130, 7, 'Ruta Endpoint', 1, 0, 'L', True)
            pdf.cell(60, 7, 'Tiempo Promedio (s)', 1, 1, 'C', True)
            pdf.set_font('Helvetica', '', 9)
            for ep in slowest:
                pdf.cell(130, 7, str(ep['path'])[:65], 1, 0, 'L')
                pdf.cell(60, 7, f"{ep['avg_time']:.3f} s", 1, 1, 'C')
        else:
            pdf.set_font('Helvetica', 'I', 9)
            pdf.cell(0, 7, 'No hay llamadas registradas a endpoints.', 0, 1, 'L')

        pdf_bytes = pdf.output()
        response = HttpResponse(pdf_bytes, content_type='application/pdf')
        response['Content-Disposition'] = 'attachment; filename="reporte_ejecutivo_performance.pdf"'
        return response

    @action(detail=False, methods=['get'], url_path='export/pptx', permission_classes=[IsAdminUser])
    def export_pptx(self, request):
        """Generates an executive PowerPoint (.pptx) presentation with Nectar Labs styling."""
        import io
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.dml.color import RGBColor
        from django.http import HttpResponse
        from datetime import datetime

        server_summary = ServerRequestLog.objects.aggregate(
            avg_response_time=Avg('response_time'),
            max_response_time=Max('response_time'),
            avg_queries=Avg('query_count'),
            total_requests=Count('id')
        )

        vitals = list(PerformanceMetric.objects.values('name').annotate(
            avg_value=Avg('value'),
            count=Count('id')
        ))

        slowest = list(ServerRequestLog.objects.values('path').annotate(
            avg_time=Avg('response_time')
        ).order_by('-avg_time')[:5])

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_slide_layout = prs.slide_layouts[6]

        DARK_BG = RGBColor(8, 12, 10)
        CARD_BG = RGBColor(16, 21, 18)
        AMBER_COLOR = RGBColor(245, 158, 11)
        EMERALD_COLOR = RGBColor(16, 185, 129)
        TEXT_WHITE = RGBColor(244, 246, 240)
        TEXT_MUTED = RGBColor(160, 170, 165)

        def add_header(slide, title_text, category="NECTAR LABS DASHBOARD"):
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            p0 = tf.paragraphs[0]
            p0.text = category
            p0.font.size = Pt(10)
            p0.font.bold = True
            p0.font.color.rgb = EMERALD_COLOR

            p1 = tf.add_paragraph()
            p1.text = title_text
            p1.font.size = Pt(22)
            p1.font.bold = True
            p1.font.color.rgb = TEXT_WHITE

        # Slide 1: Portada
        slide1 = prs.slides.add_slide(blank_slide_layout)
        background1 = slide1.background
        fill1 = background1.fill
        fill1.solid()
        fill1.fore_color.rgb = DARK_BG

        txBox1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
        tf1 = txBox1.text_frame
        tf1.word_wrap = True

        p_brand = tf1.paragraphs[0]
        p_brand.text = "NECTAR LABS ARCHITECTURE"
        p_brand.font.size = Pt(14)
        p_brand.font.bold = True
        p_brand.font.color.rgb = EMERALD_COLOR

        p_title = tf1.add_paragraph()
        p_title.text = "Informe Ejecutivo de Rendimiento & Auditoría"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = AMBER_COLOR

        p_sub = tf1.add_paragraph()
        p_sub.text = f"ms-ambar Platform | Generado: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = TEXT_MUTED

        # Slide 2: KPIs de Servidor
        slide2 = prs.slides.add_slide(blank_slide_layout)
        slide2.background.fill.solid()
        slide2.background.fill.fore_color.rgb = DARK_BG
        add_header(slide2, "Métricas Clave de Infraestructura (Backend)")

        metrics_data = [
            ("TIEMPO RESPUESTA (AVG)", f"{server_summary.get('avg_response_time') or 0:.3f} s", "Latencia promedio de solicitudes web"),
            ("TIEMPO MAXIMO (PEOR CASO)", f"{server_summary.get('max_response_time') or 0:.3f} s", "Pico de latencia registrado"),
            ("CONSULTAS DB (AVG)", f"{server_summary.get('avg_queries') or 0:.1f}", "Promedio de queries SQL por request"),
            ("TOTAL SOLICITUDES", str(server_summary.get('total_requests') or 0), "Volumen de trafico en el periodo")
        ]

        left_positions = [Inches(0.8), Inches(6.8), Inches(0.8), Inches(6.8)]
        top_positions = [Inches(1.8), Inches(1.8), Inches(4.3), Inches(4.3)]

        for idx, (label, val, desc) in enumerate(metrics_data):
            box = slide2.shapes.add_textbox(left_positions[idx], top_positions[idx], Inches(5.6), Inches(2.0))
            tf = box.text_frame
            tf.word_wrap = True

            p0 = tf.paragraphs[0]
            p0.text = label
            p0.font.size = Pt(11)
            p0.font.bold = True
            p0.font.color.rgb = AMBER_COLOR

            p1 = tf.add_paragraph()
            p1.text = val
            p1.font.size = Pt(28)
            p1.font.bold = True
            p1.font.color.rgb = TEXT_WHITE

            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(10)
            p2.font.color.rgb = TEXT_MUTED

        # Slide 3: Web Vitals & Endpoints
        slide3 = prs.slides.add_slide(blank_slide_layout)
        slide3.background.fill.solid()
        slide3.background.fill.fore_color.rgb = DARK_BG
        add_header(slide3, "Desglose de Core Web Vitals & Endpoints Críticos")

        txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True

        p_vitals = tf3.paragraphs[0]
        p_vitals.text = "• Core Web Vitals Promedio:"
        p_vitals.font.size = Pt(16)
        p_vitals.font.bold = True
        p_vitals.font.color.rgb = AMBER_COLOR

        if vitals:
            for v in vitals:
                pv = tf3.add_paragraph()
                pv.text = f"   - {v['name']}: {v['avg_value']:.2f} ms ({v['count']} muestras)"
                pv.font.size = Pt(13)
                pv.font.color.rgb = TEXT_WHITE
        else:
            pv = tf3.add_paragraph()
            pv.text = "   - Sin métricas frontend registradas aún."
            pv.font.size = Pt(13)
            pv.font.color.rgb = TEXT_MUTED

        tf3.add_paragraph()
        p_slow = tf3.add_paragraph()
        p_slow.text = "• Top Endpoints con Mayor Latencia:"
        p_slow.font.size = Pt(16)
        p_slow.font.bold = True
        p_slow.font.color.rgb = EMERALD_COLOR

        if slowest:
            for ep in slowest:
                ps = tf3.add_paragraph()
                ps.text = f"   - {ep['path']}: {ep['avg_time']:.3f} s"
                ps.font.size = Pt(13)
                ps.font.color.rgb = TEXT_WHITE
        else:
            ps = tf3.add_paragraph()
            ps.text = "   - Sin peticiones registradas."
            ps.font.size = Pt(13)
            ps.font.color.rgb = TEXT_MUTED

        stream = io.BytesIO()
        prs.save(stream)
        stream.seek(0)

        response = HttpResponse(
            stream.getvalue(),
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        response['Content-Disposition'] = 'attachment; filename="reporte_ejecutivo_performance.pptx"'
        return response


